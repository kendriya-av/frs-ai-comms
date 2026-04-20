"""
Generate a single PowerPoint slide: AI & AWS Integration — FRBSF AI Communications System
Run: python generate_pptx.py
Output: ai_aws_integration.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x16, 0x28)
GOLD    = RGBColor(0xC8, 0xA9, 0x51)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF2, 0xF8)
ORANGE  = RGBColor(0xFF, 0x99, 0x00)
DARK    = RGBColor(0x23, 0x2F, 0x3E)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
BLUE    = RGBColor(0x1B, 0x65, 0xA6)
RED     = RGBColor(0xC6, 0x28, 0x28)
GRAY    = RGBColor(0x6B, 0x7A, 0x99)
LGRAY   = RGBColor(0xD1, 0xD5, 0xDB)


def add_rect(slide, left, top, w, h, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, w, h, text, font_size=10, color=NAVY,
                 bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rich_box(slide, left, top, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    if anchor:
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(2)
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
# BACKGROUND
# ══════════════════════════════════════════════════════════════════════════════
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = LIGHT
bg.line.fill.background()

# ── Navy header bar ───────────────────────────────────────────────────────────
header = add_rect(slide, 0, 0, prs.slide_width, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(5), Inches(0.35),
             "FRBSF", 11, GOLD, True)
add_text_box(slide, Inches(0.5), Inches(0.45), Inches(8), Inches(0.5),
             "AI & AWS Integration — Communications Intelligence System", 22, WHITE, True)

# ── Subtitle bar ──────────────────────────────────────────────────────────────
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.35),
             "How Amazon Comprehend, Amazon Bedrock (Claude), and AWS infrastructure power end-to-end communications intelligence",
             11, GRAY, False)

# ══════════════════════════════════════════════════════════════════════════════
# LEFT COLUMN — AI Services
# ══════════════════════════════════════════════════════════════════════════════
col_left = Inches(0.4)
col_w = Inches(4.0)

add_text_box(slide, col_left, Inches(1.65), col_w, Inches(0.3),
             "AI SERVICES", 10, GOLD, True)

# Amazon Comprehend card
box_c = add_rect(slide, col_left, Inches(1.95), col_w, Inches(1.6), WHITE, LGRAY)
add_rich_box(slide, col_left + Inches(0.15), Inches(2.0), col_w - Inches(0.3), Inches(1.5), [
    ("🧠  Amazon Comprehend", 12, ORANGE, True),
    ("Natural Language Processing", 9, GRAY, False),
    ("", 4, GRAY, False),
    ("• Sentiment Analysis — positive / negative / neutral / mixed", 9, NAVY, False),
    ("• Text Classification — 10 inquiry categories with confidence %", 9, NAVY, False),
    ("• Key Phrase Extraction — auto-tags for routing & search", 9, NAVY, False),
    ("• Entity Detection — identifies people, orgs, dates, amounts", 9, NAVY, False),
])

# Amazon Bedrock card
box_b = add_rect(slide, col_left, Inches(3.7), col_w, Inches(2.0), WHITE, LGRAY)
add_rich_box(slide, col_left + Inches(0.15), Inches(3.75), col_w - Inches(0.3), Inches(1.9), [
    ("🤖  Amazon Bedrock (Claude 3)", 12, DARK, True),
    ("Large Language Model", 9, GRAY, False),
    ("", 4, GRAY, False),
    ("• Response Drafting — template-guided, audience-aware drafts", 9, NAVY, False),
    ("• Insights Reports — executive summaries with risk areas", 9, NAVY, False),
    ("• Risk Detection — identifies misinformation & negative trends", 9, NAVY, False),
    ("• Data Generation — synthetic test data for all categories", 9, NAVY, False),
    ("• 29 category × audience templates enforce tone & compliance", 9, NAVY, False),
])

# ══════════════════════════════════════════════════════════════════════════════
# CENTER COLUMN — Data Flow
# ══════════════════════════════════════════════════════════════════════════════
col_mid = Inches(4.65)
mid_w = Inches(4.0)

add_text_box(slide, col_mid, Inches(1.65), mid_w, Inches(0.3),
             "DATA FLOW & PIPELINE", 10, GOLD, True)

# Flow diagram card
box_f = add_rect(slide, col_mid, Inches(1.95), mid_w, Inches(3.75), WHITE, LGRAY)

flow_steps = [
    ("📥  INGEST", "Inquiries (email, web, phone)\nRSS feeds (13+ outlets)\nUploaded / generated data", NAVY),
    ("⬇", "", GOLD),
    ("🔍  CLASSIFY", "Comprehend: sentiment + category\nConfidence scoring (High/Med/Low)\nKey phrase extraction", ORANGE),
    ("⬇", "", GOLD),
    ("✍️  DRAFT", "Bedrock Claude: template-guided\nAudience-aware (public/media/stakeholder)\nHuman review required", DARK),
    ("⬇", "", GOLD),
    ("📊  ANALYZE", "Sentiment trends over time\nRisk & negative alerts\nTopic distribution & word cloud", GREEN),
    ("⬇", "", GOLD),
    ("📈  REPORT", "AI-generated executive insights\nDownload: MD / HTML / DOCX / PDF\nFull audit trail logged", BLUE),
]

y_pos = Inches(2.05)
for title, desc, color in flow_steps:
    if title == "⬇":
        add_text_box(slide, col_mid + Inches(1.7), y_pos, Inches(0.5), Inches(0.2),
                     "⬇", 12, GOLD, True, PP_ALIGN.CENTER)
        y_pos += Inches(0.22)
    else:
        add_rich_box(slide, col_mid + Inches(0.15), y_pos, mid_w - Inches(0.3), Inches(0.55), [
            (title, 10, color, True),
            (desc, 8, GRAY, False),
        ])
        y_pos += Inches(0.55)

# ══════════════════════════════════════════════════════════════════════════════
# RIGHT COLUMN — AWS Infrastructure + Metrics
# ══════════════════════════════════════════════════════════════════════════════
col_right = Inches(8.9)
right_w = Inches(4.0)

add_text_box(slide, col_right, Inches(1.65), right_w, Inches(0.3),
             "AWS INFRASTRUCTURE", 10, GOLD, True)

# Infrastructure card
box_i = add_rect(slide, col_right, Inches(1.95), right_w, Inches(1.8), WHITE, LGRAY)
add_rich_box(slide, col_right + Inches(0.15), Inches(2.0), right_w - Inches(0.3), Inches(1.7), [
    ("☁️  Deployment Architecture", 11, BLUE, True),
    ("", 4, GRAY, False),
    ("▸ App Runner — managed container hosting (HTTPS)", 9, NAVY, False),
    ("▸ ECR — Docker image registry", 9, NAVY, False),
    ("▸ CodeBuild — CI/CD (source → Docker → ECR → deploy)", 9, NAVY, False),
    ("▸ S3 — source artifact storage", 9, NAVY, False),
    ("▸ IAM — 3 roles (CodeBuild, ECR Access, Task Role)", 9, NAVY, False),
    ("▸ Region: us-east-1  |  No Docker needed locally", 9, GRAY, False),
])

# Responsible AI card
add_text_box(slide, col_right, Inches(3.85), right_w, Inches(0.3),
             "RESPONSIBLE AI", 10, GOLD, True)

box_r = add_rect(slide, col_right, Inches(4.15), right_w, Inches(1.55), WHITE, LGRAY)
add_rich_box(slide, col_right + Inches(0.15), Inches(4.2), right_w - Inches(0.3), Inches(1.45), [
    ("🛡️  Trust & Safety Guardrails", 11, RGBColor(0xC6, 0x28, 0x28), True),
    ("", 4, GRAY, False),
    ("✓ Human-in-the-loop — no auto-send path", 9, NAVY, False),
    ("✓ 9 active guardrails with live monitoring", 9, NAVY, False),
    ("✓ No policy predictions or financial advice", 9, NAVY, False),
    ("✓ Full audit trail — every AI action logged", 9, NAVY, False),
    ("✓ Confidence thresholds — low conf → mandatory review", 9, NAVY, False),
])

# ══════════════════════════════════════════════════════════════════════════════
# BOTTOM — Key Metrics Bar
# ══════════════════════════════════════════════════════════════════════════════
bar_y = Inches(5.9)
bar_h = Inches(0.9)

add_text_box(slide, Inches(0.4), Inches(5.8), Inches(3), Inches(0.25),
             "KEY CAPABILITIES", 10, GOLD, True)

metrics = [
    ("15+", "App Pages", NAVY),
    ("10", "Inquiry Categories", ORANGE),
    ("29", "Response Templates", DARK),
    ("13+", "News Outlets (RSS)", GREEN),
    ("4", "Download Formats", BLUE),
    ("9", "AI Guardrails", RED),
    ("5", "AWS AI/Infra Services", BLUE),
]

metric_w = Inches(1.7)
for i, (val, lbl, color) in enumerate(metrics):
    x = Inches(0.4) + i * (metric_w + Inches(0.1))
    box = add_rect(slide, x, bar_y, metric_w, bar_h, WHITE, LGRAY)
    add_text_box(slide, x, bar_y + Inches(0.08), metric_w, Inches(0.4),
                 val, 22, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, bar_y + Inches(0.5), metric_w, Inches(0.3),
                 lbl, 9, GRAY, False, PP_ALIGN.CENTER)

# ── Footer ────────────────────────────────────────────────────────────────────
add_text_box(slide, Inches(0.4), Inches(7.0), Inches(12), Inches(0.3),
             "Federal Reserve Bank of San Francisco  ·  External Communications  ·  AI Communications Intelligence System  ·  Powered by Amazon Comprehend + Bedrock (Claude)",
             8, GRAY, False, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
prs.save("ai_aws_integration.pptx")
print("✓ Generated: ai_aws_integration.pptx")
