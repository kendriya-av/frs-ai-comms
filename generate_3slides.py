"""
Generate 3 PowerPoint slides for FRBSF AI Communications System:
  1. Technical Excellence Highlights
  2. Responsible AI & Governance
  3. Future Enhancements
Run: python generate_3slides.py
Output: frbsf_ai_comms_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x16, 0x28)
GOLD    = RGBColor(0xC8, 0xA9, 0x51)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF2, 0xF8)
ORANGE  = RGBColor(0xFF, 0x99, 0x00)
DARK    = RGBColor(0x23, 0x2F, 0x3E)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
BLUE    = RGBColor(0x1B, 0x65, 0xA6)
RED     = RGBColor(0xC6, 0x28, 0x28)
GRAY    = RGBColor(0x6B, 0x7A, 0x99)
LGRAY   = RGBColor(0xD1, 0xD5, 0xDB)
TEAL    = RGBColor(0x00, 0x89, 0x7B)
PURPLE  = RGBColor(0x6A, 0x1B, 0x9A)


def add_rect(slide, left, top, w, h, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, w, h, text, size=10, color=NAVY,
             bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return txBox


def add_rich(slide, left, top, w, h, lines):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(2)
    return txBox


def add_slide_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()


def add_header(slide, title, subtitle):
    add_rect(slide, 0, 0, prs.slide_width, Inches(1.15), NAVY)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(3), Inches(0.3),
             "FRBSF", 11, GOLD, True)
    add_text(slide, Inches(0.5), Inches(0.42), Inches(10), Inches(0.55),
             title, 24, WHITE, True)
    add_text(slide, Inches(0.5), Inches(1.25), Inches(12), Inches(0.3),
             subtitle, 11, GRAY)


def add_footer(slide):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
             "Federal Reserve Bank of San Francisco  ·  AI Communications Intelligence System  ·  Amazon Comprehend + Bedrock (Claude)",
             8, GRAY, False, PP_ALIGN.CENTER)


def metric_box(slide, x, y, val, label, color):
    add_rect(slide, x, y, Inches(1.55), Inches(0.85), WHITE, LGRAY)
    add_text(slide, x, y + Inches(0.05), Inches(1.55), Inches(0.4),
             val, 22, color, True, PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.48), Inches(1.55), Inches(0.3),
             label, 8, GRAY, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Technical Excellence Highlights
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(s1)
add_header(s1, "Technical Excellence Highlights",
           "Production-grade architecture, real AWS services, and end-to-end automation")
add_footer(s1)

# ── Left column: Architecture ─────────────────────────────────────────────────
add_text(s1, Inches(0.4), Inches(1.65), Inches(4), Inches(0.25),
         "ARCHITECTURE & STACK", 10, GOLD, True)

add_rect(s1, Inches(0.4), Inches(1.95), Inches(4), Inches(2.2), WHITE, LGRAY)
add_rich(s1, Inches(0.55), Inches(2.0), Inches(3.7), Inches(2.1), [
    ("Full-Stack Python Application", 12, NAVY, True),
    ("", 4, GRAY, False),
    ("▸ Dash + Plotly — 15+ interactive pages, real-time charts", 9, NAVY, False),
    ("▸ Amazon Comprehend — NLP classification, sentiment, entities", 9, NAVY, False),
    ("▸ Amazon Bedrock (Claude 3) — LLM drafting, reports, risk analysis", 9, NAVY, False),
    ("▸ Pandas 3.x — data pipeline, merging, deduplication", 9, NAVY, False),
    ("▸ 13+ RSS feeds — live data from Fed, CNBC, Reuters, Bloomberg", 9, NAVY, False),
    ("▸ Zero API keys — all public data sources, no credentials needed", 9, NAVY, False),
    ("▸ Single-file app (4000+ lines) with modular service layer", 9, GRAY, False),
])

# ── Center column: AWS Infrastructure ─────────────────────────────────────────
add_text(s1, Inches(4.65), Inches(1.65), Inches(4), Inches(0.25),
         "AWS CLOUD INFRASTRUCTURE", 10, GOLD, True)

add_rect(s1, Inches(4.65), Inches(1.95), Inches(4), Inches(2.2), WHITE, LGRAY)
add_rich(s1, Inches(4.8), Inches(2.0), Inches(3.7), Inches(2.1), [
    ("Fully Deployed on AWS (us-east-1)", 12, BLUE, True),
    ("", 4, GRAY, False),
    ("▸ App Runner — managed container hosting, auto-scaling, HTTPS", 9, NAVY, False),
    ("▸ ECR — private Docker image registry", 9, NAVY, False),
    ("▸ CodeBuild — CI/CD pipeline (zip → Docker → ECR → deploy)", 9, NAVY, False),
    ("▸ S3 — source artifact storage for deployments", 9, NAVY, False),
    ("▸ IAM — 3 purpose-built roles (least privilege)", 9, NAVY, False),
    ("▸ No Docker locally — all builds run in CodeBuild", 9, NAVY, False),
    ("▸ One-command deploy: zip → S3 → CodeBuild → App Runner", 9, GRAY, False),
])

# ── Right column: Key Features ────────────────────────────────────────────────
add_text(s1, Inches(8.9), Inches(1.65), Inches(4), Inches(0.25),
         "STANDOUT FEATURES", 10, GOLD, True)

add_rect(s1, Inches(8.9), Inches(1.95), Inches(4), Inches(2.2), WHITE, LGRAY)
add_rich(s1, Inches(9.05), Inches(2.0), Inches(3.7), Inches(2.1), [
    ("What Sets This Apart", 12, GREEN, True),
    ("", 4, GRAY, False),
    ("▸ End-to-end pipeline — ingest → classify → draft → report", 9, NAVY, False),
    ("▸ 29 category × audience response templates", 9, NAVY, False),
    ("▸ Sentiment trend analysis — daily / weekly / monthly / quarterly", 9, NAVY, False),
    ("▸ ROI calculator with real AWS pricing (Bedrock, Comprehend)", 9, NAVY, False),
    ("▸ Multi-format report export — MD, HTML, DOCX, PDF", 9, NAVY, False),
    ("▸ Live data refresh — RSS feeds fetched on demand", 9, NAVY, False),
    ("▸ Synthetic data generation via Bedrock for testing", 9, GRAY, False),
])

# ── Bottom metrics row ────────────────────────────────────────────────────────
add_text(s1, Inches(0.4), Inches(4.35), Inches(4), Inches(0.25),
         "TECHNICAL METRICS", 10, GOLD, True)

metrics_1 = [
    ("4,000+", "Lines of Code", NAVY),
    ("15+", "Interactive Pages", BLUE),
    ("5", "AWS Services", ORANGE),
    ("29", "Response Templates", DARK),
    ("13+", "RSS Feed Sources", GREEN),
    ("10", "Inquiry Categories", TEAL),
    ("< $30/mo", "Total AWS Cost", RED),
]
for i, (v, l, c) in enumerate(metrics_1):
    metric_box(s1, Inches(0.4) + i * Inches(1.75), Inches(4.6), v, l, c)

# ── Bottom: Integration diagram ──────────────────────────────────────────────
add_text(s1, Inches(0.4), Inches(5.65), Inches(12), Inches(0.25),
         "INTEGRATION FLOW", 10, GOLD, True)

flow_items = [
    ("📥 Inquiries\n+ RSS Feeds", NAVY),
    ("→", GOLD),
    ("🔍 Comprehend\nClassify + Sentiment", ORANGE),
    ("→", GOLD),
    ("✍️ Bedrock Claude\nDraft + Report", DARK),
    ("→", GOLD),
    ("👤 Human Review\nApprove / Edit", GREEN),
    ("→", GOLD),
    ("📊 Dashboard\nInsights + Alerts", BLUE),
    ("→", GOLD),
    ("📋 Audit Log\nFull Traceability", GRAY),
]

x_flow = Inches(0.4)
for text, color in flow_items:
    if text == "→":
        add_text(s1, x_flow, Inches(6.05), Inches(0.4), Inches(0.5),
                 "→", 20, GOLD, True, PP_ALIGN.CENTER)
        x_flow += Inches(0.5)
    else:
        add_rect(s1, x_flow, Inches(5.9), Inches(1.8), Inches(0.85), color)
        add_text(s1, x_flow, Inches(5.98), Inches(1.8), Inches(0.7),
                 text, 9, WHITE, True, PP_ALIGN.CENTER)
        x_flow += Inches(1.9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Responsible AI & Governance
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(s2)
add_header(s2, "Responsible AI & Governance",
           "Built-in guardrails, transparency, and human oversight at every step")
add_footer(s2)

# ── Left: 9 Guardrails ───────────────────────────────────────────────────────
add_text(s2, Inches(0.4), Inches(1.65), Inches(6), Inches(0.25),
         "9 ACTIVE GUARDRAILS", 10, GOLD, True)

guardrails = [
    ("✓  Human-in-the-Loop", "No auto-send path exists. Every AI draft requires human review and approval before sending.", GREEN),
    ("✓  Closed Label Set", "Classifier uses a fixed set of 10 categories — unknown inputs map to 'other'. No hallucinated categories.", GREEN),
    ("✓  Prompt Templating", "User input is never concatenated into prompts. All Bedrock prompts use structured templates with clear boundaries.", GREEN),
    ("✓  Template-Guided Output", "29 category × audience templates enforce tone, structure, and compliance for every response.", GREEN),
    ("✓  Confidence Thresholds", "Inquiries below 60% confidence are flagged and cannot be auto-routed. Low confidence → mandatory human review.", ORANGE),
    ("✓  Full Audit Trail", "Every AI action logged with timestamp, model ID, input summary, and output summary. Complete traceability.", GREEN),
    ("✓  Bias Disclosure", "Users are warned that sentiment models may reflect training bias. Results should be interpreted with domain expertise.", ORANGE),
    ("✓  No PII Processing", "System does not store or process personally identifiable information.", GREEN),
    ("✓  No Policy Predictions", "AI is explicitly instructed to never make monetary policy predictions or provide financial advice.", RED),
]

y_g = Inches(1.95)
for title, desc, color in guardrails:
    add_rect(s2, Inches(0.4), y_g, Inches(6.0), Inches(0.52), WHITE, LGRAY)
    add_rich(s2, Inches(0.55), y_g + Inches(0.02), Inches(5.7), Inches(0.48), [
        (title, 10, color, True),
        (desc, 8, GRAY, False),
    ])
    y_g += Inches(0.56)

# ── Right top: Trust & Safety Dashboard ───────────────────────────────────────
add_text(s2, Inches(6.7), Inches(1.65), Inches(6), Inches(0.25),
         "LIVE MONITORING", 10, GOLD, True)

add_rect(s2, Inches(6.7), Inches(1.95), Inches(6.2), Inches(2.0), WHITE, LGRAY)
add_rich(s2, Inches(6.85), Inches(2.0), Inches(5.9), Inches(1.9), [
    ("🛡️  Trust & Safety Dashboard", 13, RED, True),
    ("Real-time system health — no mocks, all live data", 9, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ AWS Service Status — Bedrock, Comprehend, S3, CloudWatch, Secrets Manager", 9, NAVY, False),
    ("▸ Classification Confidence Distribution — High / Medium / Low buckets", 9, NAVY, False),
    ("▸ Auto-refreshes every 10 seconds via callback (responsive metrics)", 9, NAVY, False),
    ("▸ Guardrail Details — expandable view with implementation notes", 9, NAVY, False),
    ("▸ Audit Trail Health — action counts by type (drafts, reports, risks)", 9, NAVY, False),
])

# ── Right middle: Governance Model ────────────────────────────────────────────
add_text(s2, Inches(6.7), Inches(4.1), Inches(6), Inches(0.25),
         "GOVERNANCE MODEL", 10, GOLD, True)

add_rect(s2, Inches(6.7), Inches(4.4), Inches(6.2), Inches(2.6), WHITE, LGRAY)

gov_items = [
    ("🔒  Access Control", "IAM roles with least-privilege policies. Comprehend and Bedrock access scoped to specific API actions only.", NAVY),
    ("📋  Audit & Compliance", "Every Bedrock invocation logged: model ID, token count, input/output summary. Exportable audit log (JSON).", NAVY),
    ("📊  Explainability", "Confidence scores, key phrases, and sentiment breakdowns surfaced alongside every classification.", NAVY),
    ("⚖️  Fairness", "Sentiment models flagged for potential bias. Users instructed to apply domain expertise to AI outputs.", NAVY),
    ("🔄  Model Flexibility", "Configurable model selection (Claude 3.5 Sonnet, Sonnet, Haiku, Instant). No vendor lock-in to a single model.", NAVY),
    ("📥  Data Sovereignty", "All data processed in us-east-1. No data leaves the AWS account. No third-party APIs for core functionality.", NAVY),
]

y_gov = Inches(4.45)
for title, desc, color in gov_items:
    add_rich(s2, Inches(6.85), y_gov, Inches(5.9), Inches(0.42), [
        (title, 10, color, True),
        (desc, 8, GRAY, False),
    ])
    y_gov += Inches(0.42)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Future Enhancements
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(s3)
add_header(s3, "Future Enhancements",
           "Roadmap for scaling the platform from prototype to production")
add_footer(s3)

# ── Phase 1: Near-Term (0–3 months) ──────────────────────────────────────────
add_text(s3, Inches(0.4), Inches(1.65), Inches(4), Inches(0.25),
         "PHASE 1  ·  NEAR-TERM (0–3 MONTHS)", 10, GREEN, True)

add_rect(s3, Inches(0.4), Inches(1.95), Inches(3.9), Inches(4.5), WHITE, LGRAY)
add_rich(s3, Inches(0.55), Inches(2.0), Inches(3.6), Inches(4.4), [
    ("🚀  Quick Wins", 13, GREEN, True),
    ("", 4, GRAY, False),
    ("▸ Email Integration", 10, NAVY, True),
    ("Connect to FRBSF email system for auto-ingestion of inquiries. SES or Exchange connector.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ User Authentication", 10, NAVY, True),
    ("Add Cognito / SSO login with role-based access (officer, manager, admin).", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Response Approval Workflow", 10, NAVY, True),
    ("Multi-step approval: draft → review → approve → send. Track status per inquiry.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ CloudWatch Integration", 10, NAVY, True),
    ("Push audit logs to CloudWatch Logs. Create dashboards for AI usage metrics.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Automated Scheduling", 10, NAVY, True),
    ("EventBridge rules to auto-fetch RSS feeds and generate weekly insight reports.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Mobile-Responsive UI", 10, NAVY, True),
    ("Optimize Dash layout for tablet and mobile access by leadership.", 8, GRAY, False),
])

# ── Phase 2: Mid-Term (3–6 months) ───────────────────────────────────────────
add_text(s3, Inches(4.55), Inches(1.65), Inches(4), Inches(0.25),
         "PHASE 2  ·  MID-TERM (3–6 MONTHS)", 10, BLUE, True)

add_rect(s3, Inches(4.55), Inches(1.95), Inches(3.9), Inches(4.5), WHITE, LGRAY)
add_rich(s3, Inches(4.7), Inches(2.0), Inches(3.6), Inches(4.4), [
    ("⚡  Scale & Intelligence", 13, BLUE, True),
    ("", 4, GRAY, False),
    ("▸ Custom Comprehend Model", 10, NAVY, True),
    ("Train a custom classifier on FRBSF-specific categories for higher accuracy than keyword matching.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Knowledge Base (RAG)", 10, NAVY, True),
    ("Bedrock Knowledge Base with FRBSF policy docs, FAQs, and past responses for grounded drafting.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Multi-Language Support", 10, NAVY, True),
    ("Comprehend language detection + Bedrock translation for Spanish, Chinese, and other languages.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Predictive Analytics", 10, NAVY, True),
    ("Forecast inquiry volume spikes around FOMC meetings using historical patterns.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ DynamoDB Backend", 10, NAVY, True),
    ("Replace in-memory data with DynamoDB for persistence, multi-user support, and scalability.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ API Gateway + Lambda", 10, NAVY, True),
    ("Expose classification and drafting as REST APIs for integration with other Fed systems.", 8, GRAY, False),
])

# ── Phase 3: Long-Term (6–12 months) ─────────────────────────────────────────
add_text(s3, Inches(8.7), Inches(1.65), Inches(4.2), Inches(0.25),
         "PHASE 3  ·  LONG-TERM (6–12 MONTHS)", 10, PURPLE, True)

add_rect(s3, Inches(8.7), Inches(1.95), Inches(4.2), Inches(4.5), WHITE, LGRAY)
add_rich(s3, Inches(8.85), Inches(2.0), Inches(3.9), Inches(4.4), [
    ("🌐  Enterprise Platform", 13, PURPLE, True),
    ("", 4, GRAY, False),
    ("▸ Multi-District Deployment", 10, NAVY, True),
    ("Extend to other Federal Reserve Banks. Shared infrastructure, district-specific models.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Real-Time Streaming", 10, NAVY, True),
    ("Kinesis Data Streams for real-time social media monitoring and instant alert generation.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Fine-Tuned Foundation Model", 10, NAVY, True),
    ("Fine-tune Claude on FRBSF communication style for more authentic, on-brand responses.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Automated Response Sending", 10, NAVY, True),
    ("With sufficient guardrails and confidence, enable auto-send for high-confidence, low-risk inquiries.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ Advanced Analytics", 10, NAVY, True),
    ("Bedrock Agents for complex multi-step analysis. Correlation between sentiment and market data.", 8, GRAY, False),
    ("", 4, GRAY, False),
    ("▸ FedRAMP Compliance", 10, NAVY, True),
    ("Full security assessment for production deployment within Federal Reserve infrastructure.", 8, GRAY, False),
])

# ── Bottom: Vision statement ──────────────────────────────────────────────────
add_rect(s3, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4), NAVY)
add_text(s3, Inches(0.5), Inches(6.63), Inches(12.3), Inches(0.35),
         "Vision:  From prototype to the Federal Reserve's standard AI-powered communications platform — secure, scalable, and responsible.",
         11, WHITE, True, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
prs.save("frbsf_ai_comms_slides.pptx")
print("✓ Generated: frbsf_ai_comms_slides.pptx (3 slides)")
